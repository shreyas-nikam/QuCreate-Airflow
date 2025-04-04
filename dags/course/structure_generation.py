"""
Structure Consists of:
1. Slide_Generated: Generated slide using the content from the previous step.
2. Module Information: Markdown content for Module Information
"""
# Standard library imports
import ast
import json
import logging
import os
import re
import subprocess
from pathlib import Path

# Third-party imports
import pypandoc
from bson.objectid import ObjectId
from pptx import Presentation
from pptx.util import Inches, Pt

# Local imports
from utils.mongodb_client import AtlasClient
from utils.s3_file_manager import S3FileManager
from course.helper import convert_mermaid_diagrams_to_links


def _get_course_and_module(course_id, module_id):
    """
    Get the course and module object from the course_id and module_id
    """
    try:
        mongodb_client = AtlasClient()

        # Get the course_design object
        course = mongodb_client.find("course_design", filter={
                                     "_id": ObjectId(course_id)})
        if not course:
            return "Course not found", None

        course = course[0]
        module = next((m for m in course.get("modules", []) if m.get(
            "module_id") == ObjectId(module_id)), None)
        if not module:
            return None, "Module not found"

        return course, module
    except Exception as e:
        logging.error(f"Error in getting course and module: {e}")


def _get_resource_link(module):
    """
    Get the slide content resource link
    """
    try:
        for obj in module["pre_processed_content"]:
            if obj["resource_type"] == "Slide_Content":
                return obj["resource_link"]

        return None

    except Exception as e:
        logging.error(f"Error in getting slide content: {e}")
        return None


def _extract_content(module_name, slide_content_key):
    """
    Extract the slide content from the slide_content json file
    """
    try:
        s3_client = S3FileManager()
        response = s3_client.get_object(slide_content_key)
        slide_content = json.loads(response["Body"].read())
        
        logging.info(f"Slide content: {slide_content}")
        logging.info(type(slide_content))

        markdown = f""""""
        transcript = []

        for index, slide in enumerate(slide_content):
            if index==0:
                markdown+= f"{slide['slide_content']}\n\n"
            else:
                markdown += f"\n# {slide['slide_header']}\n\n{slide['slide_content']}\n\n"
            transcript.append(slide["speaker_notes"])
            markdown+="\n---\n"

        return markdown, transcript

    except Exception as e:
        logging.error(f"Error in extracting content: {e}")
        return None, None


def _generate_pptx(markdown, module_id, template_doc="https://qucoursify.s3.us-east-1.amazonaws.com/qucoursify/qu-create-templates/templates/style-99.pptx"):
    """
    Generate the pptx file using the markdown content, converting Mermaid diagrams to PNGs first.
    """
    try:
        output_dir = Path(f"output/{module_id}")
        output_dir.mkdir(parents=True, exist_ok=True)

        md_file_path = output_dir / "content.md"
        output_ppt_path = output_dir / "structure.pptx"
        reference_template = Path(f"dags/course/templates/{template_doc.split('/')[-1]}")

        markdown = convert_mermaid_diagrams_to_links(markdown, module_id)

        # Write markdown content to file
        with open(md_file_path, "w") as f:
            f.write(markdown)

        pypandoc.download_pandoc()
        
        # Convert Markdown to PPTX using Pandoc
        pypandoc.convert_file(
            md_file_path,
            'pptx',
            outputfile=output_ppt_path,
            extra_args=[
                '--reference-doc={}'.format(reference_template),
            ]
        )

        # Check if the PowerPoint file was created
        if not output_ppt_path.exists():
            raise Exception("Error in generating pptx")

        return str(output_ppt_path)

    except Exception as e:
        logging.error(f"Error in generating pptx: {e}")
        return None



def _format_ppt(output_ppt_path,
                speaker_notes,
                logo_path="dags/course/assets/logo.jpg",
                header='dags/course/assets/top.png'):

    Path(output_ppt_path).parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(output_ppt_path)

    del prs.slides._sldIdLst[0]

    slide_number = 0
    for slide in prs.slides:

        # if slide_number > 0:
        #     # move the title and content down by one inch
        #     top_margin = 0.7
        #     for shape in slide.shapes:
        #         if shape.has_text_frame:
        #             shape.top = Inches(top_margin)
        #             top_margin = 1.5

        # add speaker notes to all slides
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = speaker_notes[slide_number]
        except Exception as e:
            pass

        # Add logo to the bottom right corner of each slide
        left = prs.slide_width - Inches(1)
        top = prs.slide_height - Inches(0.65)
        width = Inches(1)
        slide.shapes.add_picture(logo_path, left, top, width=width)

        # Add header to the top of each slide
        # left = 0
        # top = 0
        # width = prs.slide_width
        # slide.shapes.add_picture(header, left, top, width=width)

        # Add page number to the top right corner of each slide
        # left = prs.slide_width - Inches(0.2)
        # top = Inches(0.2)
        # width = Inches(0.2)
        # slide_number_box = slide.shapes.add_textbox(
        #     left, top, width, width)
        # text_frame = slide_number_box.text_frame
        # p = text_frame.add_paragraph()
        # p.text = f"{slide_number+1}"
        # p.margin_top = 0
        # p.font.bold = False

        slide_number += 1

    prs.save(output_ppt_path)
    prs = Presentation(output_ppt_path)
    prs.save(output_ppt_path)


def _add_transcript_to_pptx(pptx_file_path, transcript):
    try:
        transcript = ast.literal_eval(transcript)
        _format_ppt(pptx_file_path, transcript)
        return pptx_file_path
    except Exception as e:
        logging.error(f"Error in adding transcript to pptx: {e}")
        return None


async def _save_file_to_s3(pptx_file_path, key):
    try:
        s3_client = S3FileManager()
        await s3_client.upload_file(pptx_file_path, key)
        resource_link = f"https://qucoursify.s3.us-east-1.amazonaws.com/{key}"
        return resource_link
    except Exception as e:
        logging.error(f"Error in saving file to s3: {e}")
        return None


def _update_slide_entry(course_id, module_id, resource_link):
    try:
        mongodb_client = AtlasClient()
        course, module = _get_course_and_module(course_id, module_id)
        s3_client = S3FileManager()

        resource = {
            "resource_link": resource_link,
            "resource_type": "Slide_Generated",
            "resource_name": "Slide Generated",
            "resource_description": "Slide Generated from the content",
            "resource_id": ObjectId()
        }

        module["pre_processed_structure"] = []

        module["pre_processed_structure"].append(resource)

        # add the other resources to the module except for the Slide_Content
        for obj in module["pre_processed_content"]:
            if obj["resource_type"] != "Slide_Content" or obj['resource_type'] != "Link":
                # copy the object to the new s3 location
                prev_location = obj["resource_link"]
                new_location = prev_location.replace(
                    "pre_processed_content", "pre_processed_structure")

                prev_location_key = prev_location.split(
                    "/")[3] + "/" + "/".join(prev_location.split("/")[4:])
                new_location_key = new_location.split(
                    "/")[3] + "/" + "/".join(new_location.split("/")[4:])

                s3_client.copy_file(prev_location_key, new_location_key)

                resource = {
                    "resource_link": new_location,
                    "resource_type": obj["resource_type"],
                    "resource_name": obj["resource_name"],
                    "resource_description": obj["resource_description"],
                    "resource_id": ObjectId()
                }
                module["pre_processed_structure"].append(resource)

        # Check: might break
        module["status"] = "Structure Review"
        course["modules"] = [module if m["module_id"] ==
                             module["module_id"] else m for m in course["modules"]]
        logging.info("Updated the course object with the new slide content")
        logging.info(f"Updated course: {course}")

        mongodb_client.update("course_design", filter={"_id": ObjectId(course_id)}, update={
            "$set": course
        })

        return True

    except Exception as e:
        logging.error(f"Error in updating entry: {e}")


async def process_structure_request(entry_id):
    """
    Steps:
    1. Get mongodb object the course id, module id in the course_design collection.
    2. Get the json file with the reviewed content from object's pre_processed_content object's resource_link of resource_type Slide_Content.
    3. Get the file from s3. Extract the sldie_header, slide_content, speaker_notes from the json file and convert the slide_header and slide_content to markdown.
    4. Generate the pptx file using the markdown file.
    5. Extract the transcript from the jso
    """
    mongodb_client = AtlasClient()
    entry = mongodb_client.find("in_structure_generation_queue", filter={
                                "_id": ObjectId(entry_id)})
    if not entry:
        return "Entry not found"
    entry = entry[0]
    course_id = entry["course_id"]
    module_id = entry["module_id"]

    course, module = _get_course_and_module(course_id, module_id)
    slide_content_link = _get_resource_link(module)
    if slide_content_link is None:
        return "Slide Content is not Generated. Please generate the slide content first."

    slide_content_key = slide_content_link.split(
        "/")[3] + "/" + "/".join(slide_content_link.split("/")[4:])
    markdown, transcript = _extract_content(slide_content_key)

    pptx_file_path = _generate_pptx(markdown, module_id)

    pptx_file_path_with_transcript = _add_transcript_to_pptx(
        pptx_file_path, transcript)
    resource_link = await _save_file_to_s3(pptx_file_path_with_transcript, resource_link)
    updated_slide = _update_slide_entry(course_id, module_id, resource_link)
    return updated_slide
