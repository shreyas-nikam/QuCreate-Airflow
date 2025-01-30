"""
Deliverables Consists of:
1. Slide_Generated: Generated slide using the content from the previous step.
2. Module Information: Markdown content for Module Information
3. Video: Video generated using the content from the previous step.
4. Assessment: Assessment generated using the content from the previous step.
"""
import ast
import os
from pydub import AudioSegment
import subprocess
from utils.s3_file_manager import S3FileManager
from utils.mongodb_client import AtlasClient
from bson.objectid import ObjectId
import logging
from pptx import Presentation
import fitz
import random
import multiprocessing
from pathlib import Path
import time
import azure.cognitiveservices.speech as speechsdk
from moviepy.editor import ImageClip, AudioFileClip, VideoFileClip, concatenate_videoclips
import json
import uuid
from utils.prompt_handler import PromptHandler
from utils.llm import LLM
from utils.retriever import Retriever
from course.structure_generation import _get_course_and_module


def _get_course_and_module(course_id, module_id):
    """
    Get the course and module object from the course_id and module_id
    """
    try:
        mongodb_client = AtlasClient()

        # Get the course_design object
        course = mongodb_client.find("course_design", filter={"_id": ObjectId(course_id)})
        if not course:
            return "Course not found", None

        course = course[0]
        module = next((m for m in course.get("modules", []) if m.get("module_id") == ObjectId(module_id)), None)
        if not module:
            return None, "Module not found"

        return course, module
    except Exception as e:
        logging.error(f"Error in getting course and module: {e}")


def _get_resources_link(course_id, module_id):
    """
    Get the slide content from the module
    """
    course, module = _get_course_and_module(course_id, module_id)
    try:
        content_link, slide_link = None, None
        for obj in module["pre_processed_structure"]:
            if obj["resource_type"] == "Slide_Generated":
                slide_link = obj["resource_link"]
        
        for obj in module["pre_processed_content"]:
            if obj["resource_type"] == "Slide_Content":
                content_link = obj["resource_link"]

        logging.info(f"Slide link: {slide_link}, Content link: {content_link}")
        
        return slide_link, content_link
    
    except Exception as e:
        logging.error(f"Error in getting slide content: {e}")
        return None
    
def _download_slide(slide_link, download_path):
    """
    Download the slide from the slide_link
    """
    try:
        logging.info(f"Downloading slide from: {slide_link}")
        s3_client = S3FileManager()
        logging.info(f"Downloading slide from: {slide_link}")
        slide_key = slide_link.split("/")[3]  + "/" + "/".join(slide_link.split("/")[4:])
        logging.info(f"Slide key: {slide_key}")
        slide_name = slide_key.split("/")[-1]
        logging.info(f"Slide name: {slide_name}")
        logging.info(f"Download path: {download_path}")
        download_path = download_path + "/" + slide_name
        s3_client.download_file(slide_key, download_path)
        return download_path
    except Exception as e:
        logging.error(f"Error in downloading slide: {e}")
        return None
    
def _get_slide_content(slide_link):
    """
    Get the slide content from the slide_content json file
    """
    try:
        logging.info(f"Getting slide content from: {slide_link}")
        s3_client = S3FileManager()
        slide_key = slide_link.split("/")[3] + "/" + "/".join(slide_link.split("/")[4:])
        logging.info(f"Slide key: {slide_key}")
        file = s3_client.get_object(slide_key)
        slide_content = file["Body"].read().decode("utf-8")
        return slide_content
    except Exception as e:
        logging.error(f"Error in getting slide content: {e}")
        return None
    
def _get_transcript_from_ppt(ppt):
    logging.info(f"Getting speaker notes from the ppt: {ppt}")
    prs = Presentation(ppt)
    speaker_notes = []
    for slide in prs.slides:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        speaker_notes.append(text_frame.text.replace("�", ""))
    return speaker_notes


def _ppt_to_pdf(input_file, output_dir):
    """
    Convert a PowerPoint file to PDF using LibreOffice.
    
    :param input_file: Path to the .pptx file
    :param output_dir: Directory where the PDF will be saved
    """
    if not os.path.exists(input_file):
        raise FileNotFoundError(f"File {input_file} does not exist")
    
    if not os.path.isdir(output_dir):
        raise NotADirectoryError(f"{output_dir} is not a valid directory")
    
    # LibreOffice command to convert to PDF
    command = [
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_file
    ]
    
    subprocess.run(command, check=True)
    print(f"Converted {input_file} to PDF and saved in {output_dir}")

def _create_images(file_path, module_id):
    # Convert the ppt to pdf
    OUTPUT_PATH = "output"
    target_folder = Path(f"{Path(OUTPUT_PATH)}/{module_id}")
    _ppt_to_pdf(file_path, target_folder)


    # Convert pdf to images
    source_file = Path(file_path.replace(".pptx", ".pdf"))
    target_folder = Path(f"{Path(OUTPUT_PATH)}/{module_id}/images")
    Path(target_folder).mkdir(parents=True, exist_ok=True)


    if os.path.exists(source_file):
        doc = fitz.open(source_file)
        for page_index in range(len(doc)):
            page = doc.load_page(page_index)
            pix = page.get_pixmap()
            output = target_folder / f"{module_id}-{page_index}.png"
            pix.save(output)
        doc.close()



def _get_audio(module_id, text, filename, voice_name):
    OUTPUT_PATH = "output"

    service_region = os.getenv("AZURE_TTS_SERVICE_REGION")
    speech_key = os.getenv("AZURE_TTS_SPEECH_KEY")


    speech_config = speechsdk.SpeechConfig(subscription=speech_key, region=service_region)
    speech_config.speech_synthesis_language = "en-US"

    voice_name_map = {
        "Female US": "en-US-AvaNeural",
        "Male US": "en-US-AndrewNeural",
        "Make UK": "en-GB-RyanNeural",
        "Female UK": "en-GB-AdaMultilingualNeural"
    }
    speech_config.speech_synthesis_voice_name = voice_name_map[voice_name]


    filepath = Path(f"{OUTPUT_PATH}/{module_id}/audio")


    # Create the folder if it does not exist
    Path(filepath).mkdir(parents=True, exist_ok=True)


    # Create the audio synthesizer
    audio_config = speechsdk.audio.AudioOutputConfig(
        filename=str(Path(f'{filepath}/{filename}')))
    speech_synthesizer = speechsdk.SpeechSynthesizer(
        speech_config=speech_config, audio_config=audio_config)


    result = speech_synthesizer.speak_text_async(text).get()


    while not os.path.exists(Path(f'{filepath}/{filename}')):
        time.sleep(random.randint(5, 10))


    # Get the audio
    while result.reason != speechsdk.ResultReason.SynthesizingAudioCompleted:


        time.sleep(random.randint(15, 20))
        result = speech_synthesizer.speak_text_async(text).get()


        while not os.path.exists(Path(f'{filepath}/{filename}')):
            time.sleep(random.randint(5, 10))


        # Checks result.
        if result.reason == speechsdk.ResultReason.SynthesizingAudioCompleted:
            break


        elif result.reason == speechsdk.ResultReason.Canceled:
            cancellation_details = result.cancellation_details
            print("Speech synthesis canceled: {}".format(
                cancellation_details.reason))
            if cancellation_details.reason == speechsdk.CancellationReason.Error:
                if cancellation_details.error_details:
                    print("Error details: {}".format(
                        cancellation_details.error_details))
            print("Did you update the subscription info?")


def _create_audio(module_id, transcript, voice_name):
    # Create the audio
    OUTPUT_PATH = "output"
    audio_folder = Path(f"{Path(OUTPUT_PATH)}/{module_id}/audio")
    Path(audio_folder).mkdir(parents=True, exist_ok=True)


    # Iterate over the speaker notes for each .txt file
    for index, speaker_notes in enumerate(transcript):

        # Get the audio
        _get_audio(module_id, speaker_notes, f"audio_{index+1}.wav", voice_name)


def _stitch_videos(module_id, updation_map):
    OUTPUT_PATH = "output"
    video_files = list(updation_map.values())
    clips = [VideoFileClip(video_file) for video_file in video_files]
    final_clip = concatenate_videoclips(clips, method="compose")
    Path(f"{OUTPUT_PATH}/{module_id}").mkdir(parents=True, exist_ok=True)
    final_clip.write_videofile(str(Path(f"{OUTPUT_PATH}/{module_id}/video.mp4")), fps=10)


def _get_questions(module_content, num_questions=10):
    trials = 5
    prompt_handler = PromptHandler()
    llm = LLM()
    while trials > 0:
        try:
            prompt = prompt_handler.get_prompt("CONTENT_TO_QUESTIONS_PROMPT")
            response = llm.get_response(prompt, inputs={"CONTENT": module_content, "NUM_QUESTIONS": num_questions})
            try:
                response = response[response.find("["):response.rfind("]") + 1]
                return json.loads(response)
            except:
                return json.loads(response)
        except Exception as e:
            logging.error(f"Error in getting questions: {e}")
            trials -= 1
            continue



def _save_questions(questions, module_id):
    OUTPUT_PATH = "output"
    print(questions)
    with open(Path(f"{OUTPUT_PATH}/{module_id}/questions.json"), "w", encoding='utf-8') as file:
        file.write(json.dumps(questions))


def _get_questions_helper(module_content, module_id, chunk_size=20000):
    overall_questions = {}
    logging.info(f"Getting questions for module: {module_id}")
    module_content_json = ast.literal_eval(module_content)
    module_content = ""
    for slide in module_content_json:
        module_content += slide["slide_content"] + "\n"
    
    try:
        questions = []
        text = module_content
        chunks = [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]
        for chunk in chunks:
            questions.extend(_get_questions(chunk, 10 if len(chunks)==1 else 5))

        for index, question in enumerate(questions):
            questions[index]["uuid"] = str(uuid.uuid4())

        _save_questions(questions, module_id)
    except Exception as e:
        logging.error(f"Error in getting questions for module {module_id}: {e}")

def _add_silence_to_audio(module_id, audio_file_name):
    OUTPUT_PATH = f"output/{module_id}/audio"
    new_audio_filename = audio_file_name.replace("_", "_with_silence_")
    audio_file_path = Path(f"{OUTPUT_PATH}/{audio_file_name}")
    new_audio_file_path = Path(f"{OUTPUT_PATH}/{new_audio_filename}")
    audio = AudioSegment.from_wav(audio_file_path)
    silence = AudioSegment.silent(duration=1000)
    audio = silence + audio
    audio.export(new_audio_file_path, format="wav")
    return new_audio_filename

def _create_video_parallel(module_id, slide_number, image):
    logging.info(f"Creating video for slide: {slide_number+1}")
    OUTPUT_PATH = "output"
    audio_folder = Path(f"{OUTPUT_PATH}/{module_id}/audio")
    video_folder = Path(f"{OUTPUT_PATH}/{module_id}/video")
    img_clip = ImageClip(image)


    # Load audio file for the current slide
    logging.info(f"Adding silence to audio for slide: {slide_number+1}")
    _add_silence_to_audio(module_id, f"audio_{slide_number+1}.wav")


    audio_file_path = Path(f"{audio_folder}/audio_with_silence_{slide_number+1}.wav")
    
    audio_file = AudioFileClip(str(audio_file_path))


    # Set the duration of the text clip to match the audio duration
    img_clip = img_clip.set_duration(audio_file.duration)


    # Create a video for this current slide
    img_clip = img_clip.set_audio(audio_file)


    img_clip = img_clip.subclip(0, audio_file.duration)
    img_clip.write_videofile(str(Path(f"{video_folder}/slide_{slide_number+1}.mp4")), 
                             fps=10, 
                             verbose=False, 
                             logger=None)
    return slide_number+1, str(Path(f"{video_folder}/slide_{slide_number+1}.mp4"))


def _create_videos(module_id):
    
    updation_map = {}
    OUTPUT_PATH = "output"
    image_folder = Path(f"{OUTPUT_PATH}/{module_id}/images")
    video_folder = Path(f"{OUTPUT_PATH}/{module_id}/video")
    Path(video_folder).mkdir(parents=True, exist_ok=True)


    # Get the sorted images
    search_dir = image_folder
    files = os.listdir(search_dir)
    files = [os.path.join(search_dir, f) for f in files]  # add path to each file
    files.sort(key=lambda x: os.path.getmtime(x))


    parallel_videos = [(module_id, i, img) for i, img in enumerate(files)]


    with multiprocessing.Pool(processes=1) as pool:
        results = pool.starmap(_create_video_parallel, parallel_videos)
        for slide_number, video_path in results:
            updation_map[slide_number] = video_path


    with open(Path(f"{OUTPUT_PATH}/{module_id}/updation_map.json"), "w", encoding='utf-8') as file:
        file.write(json.dumps(updation_map))

    logging.info("Stitching videos")
    _stitch_videos(module_id, updation_map)

def _generate_video(slide_path, module_id, voice_name):
    transcript = _get_transcript_from_ppt(slide_path)
    _create_images(slide_path, module_id)
    _create_audio(module_id, transcript, voice_name)
    _create_videos(module_id)


def _generate_assessment(module_id, slide_content):
    """
    Generate the assessment from the slide content
    """
    try:
        module_content = slide_content
        _get_questions_helper(module_content, module_id)
    except Exception as e:
        logging.error(f"Error in generating assessment: {e}")
        return None
    
    
async def _generate_chatbot(slide_content, destination, course_id, module_id):
    """
    Creates a retriever object for the course
    
    Args:
    file_name: Path to the folder containing the course content
    destination: Path to the folder where the retriever object will be saved
    """

    logging.info("Creating the retriever object")
    
    slide_content_json = ast.literal_eval(slide_content)
    slide_content = ""
    for slide in slide_content_json:
        slide_content += slide["slide_content"] + "\n"

    
    s3_file_manager = S3FileManager()
    with open(f"{destination}/slide_content.txt", "w") as file:
        file.write(slide_content)
    
    key = f"qu-course-design/{course_id}/{module_id}/pre_processed_deliverables/retriever.txt"
    await s3_file_manager.upload_file(f"{destination}/slide_content.txt", key)

    chabot_link = "https://qucoursify.s3.us-east-1.amazonaws.com/"+key
    return chabot_link


async def upload_files(course_id, video_path, assessment_path, chatbot_path, module_id, has_assessment, has_chatbot):
    """
    Upload the video, assessment file, chatbot to s3
    """
    try:
        logging.info(f"Uploading video to s3 for module: {module_id}")
        s3_client = S3FileManager()
        key = f"qu-course-design/{course_id}/{module_id}/pre_processed_deliverables/"
        video_key = key + "video.mp4"
        await s3_client.upload_video(video_path, video_key)
        video_link = "https://qucoursify.s3.us-east-1.amazonaws.com/"+video_key

        if has_assessment:
            logging.info(f"Uploading assessment to s3 for module: {module_id}")
            assessment_key = key + f"{module_id}_assessment.json"
            await s3_client.upload_file(assessment_path, assessment_key)
            assessment_link = "https://qucoursify.s3.us-east-1.amazonaws.com/"+assessment_key
        else:
            assessment_link = ""

        if has_chatbot:
            logging.info(f"Uploading chatbot to s3 for module: {module_id}")
            chatbot_link = chatbot_path

        else:
            chatbot_link = ""

        return video_link, assessment_link, chatbot_link
    
    except Exception as e:
        logging.error(f"Error in uploading files: {e}")
        return None, None, None
    
def update_module_with_deliverables(course_id, module_id, video_link, assessment_link, chatbot_link, has_chatbot, has_assessment, slide_link):
    """
    Update the module with the video, assessment, chatbot links in pre_processed_deliverables
    """
    course, module = _get_course_and_module(course_id, module_id)
    try:
        module["pre_processed_deliverables"] = [
            {
                "resource_type": "Video",
                "resource_link": video_link,
                "resource_name": f"{module_id}.mp4",
                "resource_description": "Video generated from the slide content",
                "resource_id": ObjectId()
            },
        ]

        for resource in module['pre_processed_structure']:
            if resource["resource_type"] == "Note":
                
                prev_location = resource["resource_link"]
                prev_location_key = prev_location.split("/")[3] + "/" + "/".join(prev_location.split("/")[4:])
                new_location = prev_location.replace("pre_processed_structure", "pre_processed_deliverables")
                new_location_key = new_location.split("/")[3] + "/" + "/".join(new_location.split("/")[4:])
                s3_client = S3FileManager()
                logging.info(f"Copying file from {prev_location_key} to {new_location_key}")
                s3_client.copy_file(prev_location_key, new_location_key)
                new_location_link = "https://qucoursify.s3.us-east-1.amazonaws.com/"+new_location_key

                module["pre_processed_deliverables"].append({
                    "resource_type": "Note",
                    "resource_link": new_location_link,
                    "resource_name": f"{module_id}_module_info.md",
                    "resource_description": "Module Information",
                    "resource_id": ObjectId()
                })

        module["pre_processed_deliverables"].append({
            "resource_type": "Slide_Generated",
            "resource_link": slide_link,
            "resource_name": f"{module_id}.pptx",
            "resource_description": "Slides generated from the slide content",
            "resource_id": ObjectId()
        })

        if has_assessment:
            module["pre_processed_deliverables"].append({
                "resource_type": "Assessment",
                "resource_link": assessment_link,
                "resource_name": f"{module_id}_assessment.json",
                "resource_description": "Assessment generated from the slide content",
                "resource_id": ObjectId()
            })

        if has_chatbot:
            module["pre_processed_deliverables"].append({
                "resource_type": "Chatbot",
                "resource_link": chatbot_link,
                "resource_name": f"Chatbot.txt",
                "resource_description": "Chatbot text for creating the retriever",
                "resource_id": ObjectId()
            })
        

        module["status"] = "Deliverables Review"
        course['modules'] = [module if m["module_id"] == module["module_id"] else m for m in course['modules']]
        mongodb_client = AtlasClient()
        mongodb_client.update("course_design", filter={"_id": ObjectId(course_id)}, update={"$set": course})
        
    except Exception as e:
        logging.error(f"Error in updating module with deliverables: {e}")

