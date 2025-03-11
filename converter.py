import os
import re
import argparse
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
import tempfile
from PIL import Image
import yaml
import subprocess
import cairosvg
import io

class MarkdownToPPT:
    def __init__(self, template_path=None):
        """Initialize the converter with an optional template."""
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
        
        self.slide_layouts = {
            'title': 0,  # Title Slide
            'section': 1,  # Section Header
            'content': 2,  # Title and Content
            'two_content': 3,  # Two Content
            'title_only': 5,  # Title Only
            'blank': 6,  # Blank
        }
        
        # Default theme colors
        self.theme = {
            'primary': RGBColor(31, 73, 125),
            'secondary': RGBColor(79, 129, 189),
            'background': RGBColor(255, 255, 255),
            'text': RGBColor(0, 0, 0),
            'accent1': RGBColor(192, 80, 77),
            'accent2': RGBColor(155, 187, 89),
        }
        
        # Default font settings
        self.fonts = {
            'title': {
                'name': 'Calibri',
                'size': Pt(44),
                'bold': True,
                'color': self.theme['primary']
            },
            'subtitle': {
                'name': 'Calibri',
                'size': Pt(32),
                'bold': False,
                'color': self.theme['secondary']
            },
            'heading': {
                'name': 'Calibri',
                'size': Pt(28),
                'bold': True,
                'color': self.theme['primary']
            },
            'body': {
                'name': 'Calibri',
                'size': Pt(18),
                'bold': False,
                'color': self.theme['text']
            }
        }
        
        # Temporary directory for generated images
        self.temp_dir = tempfile.mkdtemp()
    
    def load_config(self, config_path):
        """Load configuration from YAML file."""
        if not os.path.exists(config_path):
            return
        
        with open(config_path, 'r') as file:
            config = yaml.safe_load(file)
        
        # Update theme colors
        if 'theme' in config:
            for key, color_hex in config['theme'].items():
                if key in self.theme:
                    r, g, b = self._hex_to_rgb(color_hex)
                    self.theme[key] = RGBColor(r, g, b)
        
        # Update font settings
        if 'fonts' in config:
            for font_type, settings in config['fonts'].items():
                if font_type in self.fonts:
                    for setting, value in settings.items():
                        if setting == 'size':
                            self.fonts[font_type][setting] = Pt(value)
                        elif setting == 'color':
                            r, g, b = self._hex_to_rgb(value)
                            self.fonts[font_type][setting] = RGBColor(r, g, b)
                        else:
                            self.fonts[font_type][setting] = value
        
        # Update slide layouts
        if 'layouts' in config:
            for layout_name, layout_index in config['layouts'].items():
                self.slide_layouts[layout_name] = layout_index
    
    def _hex_to_rgb(self, hex_color):
        """Convert hex color to RGB tuple."""
        hex_color = hex_color.lstrip('#')
        return int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    
    def create_template(self, output_path):
        """Create a template configuration file."""
        template_config = {
            'theme': {
                'primary': '#1F497D',
                'secondary': '#4F81BD',
                'background': '#FFFFFF',
                'text': '#000000',
                'accent1': '#C0504D',
                'accent2': '#9BBB59'
            },
            'fonts': {
                'title': {
                    'name': 'Calibri',
                    'size': 44,
                    'bold': True,
                    'color': '#1F497D'
                },
                'subtitle': {
                    'name': 'Calibri',
                    'size': 32,
                    'bold': False,
                    'color': '#4F81BD'
                },
                'heading': {
                    'name': 'Calibri',
                    'size': 28,
                    'bold': True,
                    'color': '#1F497D'
                },
                'body': {
                    'name': 'Calibri',
                    'size': 18,
                    'bold': False,
                    'color': '#000000'
                }
            },
            'layouts': {
                'title': 0,
                'section': 1,
                'content': 2,
                'two_content': 3,
                'title_only': 5,
                'blank': 6
            }
        }
        
        with open(output_path, 'w') as file:
            yaml.dump(template_config, file, default_flow_style=False)
        
        print(f"Template configuration saved to {output_path}")
    
    def convert(self, markdown_file, output_file):
        """Convert markdown file to PowerPoint presentation."""
        with open(markdown_file, 'r', encoding='utf-8') as file:
            md_content = file.read()
        
        # Process Mermaid diagrams
        md_content = self._process_mermaid_blocks(md_content)
        
        # Convert markdown to HTML
        html = markdown.markdown(md_content, extensions=['extra', 'toc', 'sane_lists', 'codehilite'])
        
        # Parse HTML
        soup = BeautifulSoup(html, 'html.parser')
        
        # Process slides
        self._process_slides(soup)
        
        # Save presentation
        self.prs.save(output_file)
        print(f"Presentation saved to {output_file}")
    
    def _process_mermaid_blocks(self, md_content):
        """Process Mermaid code blocks and convert them to images."""
        mermaid_pattern = re.compile(r'```mermaid\n(.*?)\n```', re.DOTALL)
        
        def replace_mermaid(match):
            mermaid_code = match.group(1)
            image_path = self._generate_mermaid_image(mermaid_code)
            return f"![Mermaid Diagram]({image_path})"
        
        return mermaid_pattern.sub(replace_mermaid, md_content)
    
    def _generate_mermaid_image(self, mermaid_code):
        """Generate an image from Mermaid code."""
        # Save Mermaid code to a temporary file
        mermaid_file = os.path.join(self.temp_dir, f"mermaid_{hash(mermaid_code)}.mmd")
        with open(mermaid_file, 'w') as file:
            file.write(mermaid_code)
        
        # Generate SVG using Mermaid CLI (requires Node.js and mermaid-cli)
        svg_file = os.path.join(self.temp_dir, f"mermaid_{hash(mermaid_code)}.svg")
        png_file = os.path.join(self.temp_dir, f"mermaid_{hash(mermaid_code)}.png")
        
        try:
            subprocess.run(
                ["mmdc", "-i", mermaid_file, "-o", svg_file, "-b", "transparent"],
                check=True
            )
            
            # Convert SVG to PNG
            cairosvg.svg2png(url=svg_file, write_to=png_file)
            return png_file
        except Exception as e:
            print(f"Failed to generate Mermaid diagram: {e}")
            # Create a simple placeholder image
            placeholder = Image.new('RGB', (800, 400), color=(255, 255, 255))
            placeholder.save(png_file)
            return png_file
    
    def _process_slides(self, soup):
        """Process HTML and create slides."""
        # Find all headings that should start a new slide
        slide_headers = soup.find_all(['h1', 'h2'])
        
        if not slide_headers:
            # If no headers found, create a single slide with all content
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.slide_layouts['content']])
            title = slide.shapes.title
            title.text = "Presentation"
            content = soup.get_text()
            text_box = slide.shapes.placeholders[1]
            text_frame = text_box.text_frame
            p = text_frame.add_paragraph()
            p.text = content
            self._format_paragraph(p, 'body')
            return
        
        # Process each section as a slide
        for i, header in enumerate(slide_headers):
            # Determine slide layout based on header level
            layout_key = 'title' if header.name == 'h1' and i == 0 else 'section' if header.name == 'h1' else 'content'
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.slide_layouts[layout_key]])
            
            # Set slide title
            title = slide.shapes.title
            title.text = header.get_text()
            
            # Get content until next header or end of document
            content_elements = []
            next_sibling = header.next_sibling
            
            while next_sibling and next_sibling.name not in ['h1', 'h2']:
                if next_sibling.name:  # Skip NavigableString objects
                    content_elements.append(next_sibling)
                next_sibling = next_sibling.next_sibling
                if not next_sibling:
                    break
            
            # Add content to slide
            if layout_key == 'title':
                # Look for a subtitle (h2 right after h1)
                if content_elements and content_elements[0].name == 'h2':
                    subtitle = slide.placeholders[1]
                    subtitle.text = content_elements[0].get_text()
                    content_elements = content_elements[1:]
            
            # Process remaining content
            if content_elements:
                if layout_key != 'title':  # Title slides don't have a content placeholder
                    try:
                        # Try to get the content placeholder (index may vary)
                        content_placeholder = None
                        for shape in slide.placeholders:
                            if shape.placeholder_format.idx in [1, 2]:  # Common indices for content placeholders
                                content_placeholder = shape
                                break
                        
                        if content_placeholder:
                            self._add_content_to_placeholder(content_placeholder, content_elements)
                        else:
                            # Create a text box if no placeholder found
                            self._add_content_to_textbox(slide, content_elements)
                    except Exception as e:
                        print(f"Error adding content to slide: {e}")
                        # Fallback to text box
                        self._add_content_to_textbox(slide, content_elements)
                else:
                    # For title slides, add a text box below for any additional content
                    self._add_content_to_textbox(slide, content_elements)
    
    def _add_content_to_placeholder(self, placeholder, content_elements):
        """Add content to a slide placeholder."""
        text_frame = placeholder.text_frame
        text_frame.clear()  # Clear existing content
        
        for element in content_elements:
            if element.name in ['p', 'div']:
                p = text_frame.add_paragraph()
                p.text = element.get_text()
                self._format_paragraph(p, 'body')
            elif element.name in ['h3', 'h4', 'h5', 'h6']:
                p = text_frame.add_paragraph()
                p.text = element.get_text()
                self._format_paragraph(p, 'heading')
            elif element.name == 'ul':
                self._add_list_to_frame(text_frame, element, is_numbered=False)
            elif element.name == 'ol':
                self._add_list_to_frame(text_frame, element, is_numbered=True)
            elif element.name == 'img':
                self._add_image_after_placeholder(placeholder, element)
            elif element.name == 'pre':
                p = text_frame.add_paragraph()
                p.text = element.get_text()
                p.font.name = 'Courier New'
                p.font.size = Pt(14)
    
    def _add_content_to_textbox(self, slide, content_elements):
        """Add content to a new text box on the slide."""
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        for element in content_elements:
            if element.name in ['p', 'div']:
                p = text_frame.add_paragraph()
                p.text = element.get_text()
                self._format_paragraph(p, 'body')
            elif element.name in ['h3', 'h4', 'h5', 'h6']:
                p = text_frame.add_paragraph()
                p.text = element.get_text()
                self._format_paragraph(p, 'heading')
            elif element.name == 'ul':
                self._add_list_to_frame(text_frame, element, is_numbered=False)
            elif element.name == 'ol':
                self._add_list_to_frame(text_frame, element, is_numbered=True)
            elif element.name == 'img':
                self._add_image_to_slide(slide, element)
            elif element.name == 'pre':
                p = text_frame.add_paragraph()
                p.text = element.get_text()
                p.font.name = 'Courier New'
                p.font.size = Pt(14)
    
    def _add_list_to_frame(self, text_frame, list_element, is_numbered=False):
        """Add a list to a text frame."""
        items = list_element.find_all('li')
        for i, item in enumerate(items):
            p = text_frame.add_paragraph()
            prefix = f"{i+1}. " if is_numbered else "• "
            p.text = prefix + item.get_text()
            self._format_paragraph(p, 'body')
            p.level = 0  # Set bullet level
    
    def _add_image_after_placeholder(self, placeholder, img_element):
        """Add an image after a placeholder on the slide."""
        slide = placeholder.part.slide
        self._add_image_to_slide(slide, img_element, placeholder.top + placeholder.height)
    
    def _add_image_to_slide(self, slide, img_element, top_position=None):
        """Add an image to a slide."""
        src = img_element.get('src', '')
        alt = img_element.get('alt', 'Image')
        
        if not src:
            return
        
        image_path = src
        if src.startswith(('http://', 'https://')):
            # Download the image
            try:
                response = requests.get(src, stream=True)
                if response.status_code == 200:
                    image_path = os.path.join(self.temp_dir, f"img_{hash(src)}.png")
                    with open(image_path, 'wb') as file:
                        for chunk in response:
                            file.write(chunk)
                else:
                    print(f"Failed to download image: {src}")
                    return
            except Exception as e:
                print(f"Error downloading image: {e}")
                return
        
        try:
            # Default image position and size
            left = Inches(1)
            top = Inches(3) if top_position is None else top_position + Inches(0.5)
            width = Inches(8)
            
            # Get original image dimensions to maintain aspect ratio
            with Image.open(image_path) as img:
                orig_width, orig_height = img.size
                ratio = orig_height / orig_width
                height = width * ratio
            
            # Add image to slide
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception as e:
            print(f"Error adding image to slide: {e}")
    
    def _format_paragraph(self, paragraph, style_type):
        """Apply formatting to a paragraph based on style type."""
        if style_type not in self.fonts:
            style_type = 'body'  # Default to body style
        
        font_settings = self.fonts[style_type]
        paragraph.font.name = font_settings['name']
        paragraph.font.size = font_settings['size']
        paragraph.font.bold = font_settings['bold']
        paragraph.font.color.rgb = font_settings['color']
        
        # Set paragraph alignment
        if style_type in ['title', 'heading']:
            paragraph.alignment = PP_ALIGN.CENTER
        else:
            paragraph.alignment = PP_ALIGN.LEFT

def main():
    parser = argparse.ArgumentParser(description='Convert Markdown to PowerPoint')
    parser.add_argument('markdown_file', help='Path to markdown file')
    parser.add_argument('output_file', help='Path for output PowerPoint file')
    parser.add_argument('--template', help='Path to template PowerPoint file (optional)')
    parser.add_argument('--config', help='Path to configuration YAML file (optional)')
    parser.add_argument('--create-template', help='Create a template configuration file at the specified path')
    
    args = parser.parse_args()
    
    if args.create_template:
        converter = MarkdownToPPT()
        converter.create_template(args.create_template)
        return
    
    converter = MarkdownToPPT(args.template)
    
    if args.config:
        converter.load_config(args.config)
    
    converter.convert(args.markdown_file, args.output_file)

if __name__ == "__main__":
    main()